#!/usr/bin/env python3
"""
Manus AI 슬라이드 생성 헬퍼 스크립트

04_SlidePrompt/*.md 프롬프트 파일을 Manus AI에 일괄 제출하여
Nano Banana Pro로 슬라이드를 생성하고 PPTX를 다운로드합니다.

사용법:
  python .agent/scripts/manus_slide.py [프로젝트폴더]
  python .agent/scripts/manus_slide.py                           # 자동 탐색
  python .agent/scripts/manus_slide.py /path/to/project          # 명시적 지정
  python .agent/scripts/manus_slide.py --resume                  # 이전 완료 파일 자동 스킵
  python .agent/scripts/manus_slide.py --quiet                   # 사일런스 모드
  python .agent/scripts/manus_slide.py --verbose                 # 디버그 출력
"""

import argparse
import atexit
import json
import logging
import os
import shutil
import signal
import sys
import tempfile
import time
from datetime import datetime
from pathlib import Path

try:
    import requests
    from requests.adapters import HTTPAdapter
except ImportError:
    print("[ERROR] requests 패키지가 필요합니다: pip install requests")
    sys.exit(1)

try:
    from urllib3.util.retry import Retry
except ImportError:
    Retry = None


# ─── 로거 ────────────────────────────────────────────────────
log = logging.getLogger("manus_slide")


# ─── 설정 ───────────────────────────────────────────────────
MANUS_API_BASE = "https://api.manus.ai/v1"
POLL_INTERVAL = 30  # 폴링 간격 (초)
MAX_WAIT_TIME = 1800  # 최대 대기 시간 (30분)
MAX_CONCURRENT = 5  # 동시 제출 최대 수
PROMPT_GLOB = "*슬라이드 생성 프롬프트.md"
CHUNK_THRESHOLD = 1000  # 이 줄 수 이상이면 세션 단위 분할
CHUNK_MAX_SLIDES = 35  # 이 슬라이드 수 이상이면 분할

# ─── 안정성 설정 ─────────────────────────────────────────────
RETRY_TOTAL = 3  # 최대 재시도 횟수
RETRY_BACKOFF = 1  # 재시도 대기 (1s, 2s, 4s)
RETRY_STATUS_CODES = [429, 500, 502, 503, 504]
MIN_PPTX_SIZE = 10240  # 최소 유효 PPTX 크기 (10KB)
MIN_DISK_MB = 500  # 최소 디스크 여유 공간 (MB)
LOCK_FILENAME = ".manus_slide.lock"

# ─── Graceful Shutdown ───────────────────────────────────────
_shutdown_requested = False
_current_task_log = []
_current_output_dir = None


def _handle_signal(signum, frame):
    """SIGINT/SIGTERM 시 현재 작업 완료 후 종료"""
    global _shutdown_requested
    _shutdown_requested = True
    sig_name = "SIGINT" if signum == signal.SIGINT else "SIGTERM"
    log.warning("시그널 %s 수신 — 현재 작업 완료 후 안전 종료합니다", sig_name)


def _save_checkpoint():
    """중단 시 현재까지의 task_log를 저장"""
    if _current_output_dir and _current_task_log:
        log_path = Path(_current_output_dir) / "manus_task_log.json"
        try:
            log_path.write_text(
                json.dumps(_current_task_log, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
            log.info(
                "체크포인트 저장: %s (%d건)", log_path.name, len(_current_task_log)
            )
        except OSError as e:
            log.error("체크포인트 저장 실패: %s", e)


# ─── 헤더/풋터 금지 프리픽스 ─────────────────────────────────
# Manus AI에 프롬프트 전송 시 최상단에 삽입하여 디자인 규칙을 강제합니다.
SLIDE_GENERATION_PREFIX = """[최우선 규칙 — 모든 슬라이드에 예외 없이 적용]

■ 슬라이드 생성 방식 (필수)
반드시 Nano Banana Pro 이미지 생성을 사용하여 각 슬라이드를 고품질 AI 이미지로 생성하세요.
HTML 기반 텍스트 슬라이드가 아닌, 각 슬라이드가 AI가 생성한 이미지여야 합니다.
확인이나 진행 여부를 묻지 말고 바로 슬라이드 이미지 생성을 시작하세요.

■ 파일 내보내기 (필수)
모든 슬라이드 이미지 생성이 완료되면 반드시 다음을 수행하세요:
1. 생성된 모든 슬라이드 이미지를 하나의 PPTX 파일로 조립하세요
2. 각 슬라이드의 발표자 노트(Speaker Notes)를 PPTX에 포함하세요
3. 완성된 PPTX 파일을 다운로드 가능한 형태로 출력하세요
4. 발표자 노트를 별도 마크다운 파일(slide_notes.md)로도 저장하세요
PPTX 파일 출력을 절대 생략하지 마세요.

■ 디자인 규칙
1. 슬라이드에 상단 헤더 바(header bar)를 절대 배치하지 마세요.
2. 슬라이드에 하단 풋터 바(footer bar)를 절대 배치하지 마세요.
3. 페이지 번호(page number)를 표시하지 마세요.
4. 로고 영역을 배치하지 마세요.
5. 콘텐츠가 슬라이드 전체 영역(full bleed)을 활용해야 합니다.
6. 슬라이드 상단이나 하단에 반복되는 장식 요소(줄, 바, 밴드)를 넣지 마세요.

■ 발표자 노트(Speaker Notes) 규칙 — 필수
모든 슬라이드에 발표자 노트(Speaker Notes)를 반드시 작성하세요. 발표자 노트는 프로그래밍을 처음 가르치는 초보 강사가 슬라이드 화면만 보면서 막힘 없이 발표할 수 있도록, 아래 항목을 포함하여 충분히 자세하게 작성합니다:

1. **도입 멘트**: 이 슬라이드를 처음 보여줄 때 강사가 말할 오프닝 한두 문장
2. **핵심 설명**: 슬라이드의 핵심 내용을 학습자에게 어떻게 설명할지 구어체로 풀어쓴 3~5문장
3. **비유/예시 안내**: 슬라이드에 비유가 있으면 "이렇게 설명하세요"라는 구체적 화법 제시
4. **학습자 반응 예상**: "이 부분에서 학습자들이 자주 하는 질문은..." 또는 "여기서 막히는 분이 있으면..."
5. **시간 안배**: "이 슬라이드는 약 2분, 실습 포함 시 5분 정도 소요됩니다"
6. **전환 멘트**: 다음 슬라이드로 넘어갈 때의 연결 문장

발표자 노트 분량 기준: 슬라이드당 최소 3문장, 실습/코드 슬라이드는 최소 5문장.
발표자 노트가 없는 슬라이드가 있으면 안 됩니다.

---

"""


# ═══════════════════════════════════════════════════════════════
# 유틸리티 함수
# ═══════════════════════════════════════════════════════════════


def _setup_logging(quiet=False, verbose=False):
    """로깅 설정 (M1+M2: logging 모듈 + --quiet/--verbose)"""
    if verbose:
        level = logging.DEBUG
    elif quiet:
        level = logging.WARNING
    else:
        level = logging.INFO

    fmt = "%(asctime)s [%(levelname)s] %(message)s"
    datefmt = "%H:%M:%S"
    logging.basicConfig(level=level, format=fmt, datefmt=datefmt, stream=sys.stdout)
    log.setLevel(level)


def _create_session(api_key):
    """requests.Session 생성 + 자동 재시도 설정 (C7+C1+C5)"""
    session = requests.Session()
    session.headers.update(
        {
            "API_KEY": api_key,
            "Content-Type": "application/json",
        }
    )

    if Retry is not None:
        retry = Retry(
            total=RETRY_TOTAL,
            backoff_factor=RETRY_BACKOFF,
            status_forcelist=RETRY_STATUS_CODES,
            allowed_methods=["GET", "POST"],
            respect_retry_after_header=True,
            raise_on_status=False,
        )
        adapter = HTTPAdapter(max_retries=retry, pool_maxsize=10)
        session.mount("https://", adapter)
        log.debug(
            "requests.Session 생성 (retry=%d, backoff=%ds, status=%s)",
            RETRY_TOTAL,
            RETRY_BACKOFF,
            RETRY_STATUS_CODES,
        )
    else:
        log.warning("urllib3.Retry 사용 불가 — 재시도 없이 실행됩니다")

    return session


def _preflight_check(session):
    """배치 시작 전 API 연결 + 인증 확인 (M3)"""
    try:
        resp = session.get(
            f"{MANUS_API_BASE}/tasks",
            params={"limit": 1},
            timeout=(10, 30),
        )
        if resp.status_code == 401:
            log.error("API 인증 실패 — MANUS_API_KEY를 확인하세요")
            return False
        if resp.status_code == 403:
            log.error("API 접근 거부 — Manus Pro/Team 플랜이 필요합니다")
            return False
        if resp.status_code >= 500:
            log.warning(
                "Manus API 서버 오류 (%d) — 서비스 상태를 확인하세요", resp.status_code
            )
            return False
        log.info("Manus API 연결 확인됨 (status=%d)", resp.status_code)
        return True
    except requests.RequestException as e:
        log.error("Manus API 연결 실패: %s", e)
        return False


def _check_disk_space(output_dir, min_mb=MIN_DISK_MB):
    """디스크 여유 공간 확인 (M6)"""
    try:
        usage = shutil.disk_usage(str(output_dir))
        free_mb = usage.free // (1024 * 1024)
        if free_mb < min_mb:
            log.error("디스크 공간 부족: %dMB 남음 (최소 %dMB 필요)", free_mb, min_mb)
            return False
        log.debug("디스크 여유 공간: %dMB", free_mb)
        return True
    except OSError as e:
        log.warning("디스크 공간 확인 실패: %s", e)
        return True  # 확인 실패 시 진행 허용


def _acquire_lock(output_dir):
    """중복 실행 방지 Lock File (M8)"""
    lock_path = Path(output_dir) / LOCK_FILENAME
    if lock_path.exists():
        try:
            lock_data = json.loads(lock_path.read_text(encoding="utf-8"))
            pid = lock_data.get("pid", 0)
            # 프로세스가 실제로 살아있는지 확인
            try:
                os.kill(pid, 0)
                log.error(
                    "다른 프로세스(PID %d)가 이미 실행 중입니다. 강제 해제: %s 삭제",
                    pid,
                    lock_path,
                )
                return False
            except OSError:
                log.warning("이전 실행의 잔여 Lock File 제거 (PID %d 종료됨)", pid)
                lock_path.unlink(missing_ok=True)
        except (json.JSONDecodeError, KeyError):
            lock_path.unlink(missing_ok=True)

    lock_data = {
        "pid": os.getpid(),
        "started_at": datetime.now().isoformat(),
    }
    lock_path.write_text(json.dumps(lock_data), encoding="utf-8")
    log.debug("Lock 획득: PID %d", os.getpid())
    return True


def _release_lock(output_dir):
    """Lock File 해제"""
    lock_path = Path(output_dir) / LOCK_FILENAME
    lock_path.unlink(missing_ok=True)
    log.debug("Lock 해제")


def _load_completed_files(output_dir):
    """이전 실행에서 완료된 파일 목록 로드 (C2: --resume)"""
    log_path = Path(output_dir) / "manus_task_log.json"
    if not log_path.exists():
        return set()
    try:
        log_data = json.loads(log_path.read_text(encoding="utf-8"))
        completed = {
            r["file"]
            for r in log_data
            if r.get("status") in ("completed", "completed_no_file")
        }
        if completed:
            log.info("이전 완료 파일 %d개 발견 — 자동 스킵합니다", len(completed))
        return completed
    except (json.JSONDecodeError, KeyError, OSError) as e:
        log.warning("이전 로그 파일 읽기 실패: %s", e)
        return set()


def _cleanup_chunk_dirs(output_dir):
    """임시 청크 디렉토리 정리 (M7)"""
    for d in Path(output_dir).glob(".chunks_*"):
        if d.is_dir():
            try:
                shutil.rmtree(d)
                log.debug("임시 청크 디렉토리 삭제: %s", d.name)
            except OSError as e:
                log.warning("임시 디렉토리 삭제 실패 (%s): %s", d.name, e)


def get_api_key():
    """환경변수 또는 .agent/.env에서 MANUS_API_KEY 로드"""
    key = os.environ.get("MANUS_API_KEY")
    if key:
        return key

    env_paths = [
        Path.cwd() / ".claude" / ".env",
        Path(__file__).parent.parent / ".claude" / ".env",
    ]
    for env_path in env_paths:
        if env_path.exists():
            for line in env_path.read_text(encoding="utf-8").splitlines():
                line = line.strip()
                if line.startswith("MANUS_API_KEY="):
                    val = line.split("=", 1)[1].strip().strip('"').strip("'")
                    if val and not val.startswith("your_"):
                        return val
    return None


def find_project_folder(arg_path=None):
    """프로젝트 폴더 자동 탐색"""
    if arg_path:
        p = Path(arg_path)
        if (p / "04_SlidePrompt").is_dir():
            return p
        # 날짜 프리픽스 프로젝트 폴더 탐색
        for d in sorted(p.glob("20*"), reverse=True):
            if (d / "04_SlidePrompt").is_dir():
                return d

    cwd = Path.cwd()
    # 현재 디렉토리에 04_SlidePrompt가 있으면 바로 사용
    if (cwd / "04_SlidePrompt").is_dir():
        return cwd
    # 하위 프로젝트 폴더 탐색
    for d in sorted(cwd.glob("20*"), reverse=True):
        if (d / "04_SlidePrompt").is_dir():
            return d

    return None


def discover_prompts(project_dir):
    """04_SlidePrompt/ 내 프롬프트 파일 목록 반환"""
    prompt_dir = project_dir / "04_SlidePrompt"
    files = sorted(prompt_dir.glob(PROMPT_GLOB))
    return files


# ═══════════════════════════════════════════════════════════════
# API 호출 함수 (session 기반)
# ═══════════════════════════════════════════════════════════════


def create_task(session, prompt_content, filename):
    """Manus AI에 슬라이드 생성 Task 생성 (헤더/풋터 금지 프리픽스 자동 삽입)"""
    # 프롬프트 최상단에 규칙 프리픽스, 최하단에 export 지시 삽입
    export_suffix = "\n\n---\n\n[최종 출력 지시]\n위 교안의 모든 슬라이드를 Nano Banana Pro 이미지로 생성한 후, 반드시 하나의 PPTX 파일로 조립하여 다운로드 가능한 파일로 출력하세요. 발표자 노트도 PPTX 안에 포함하고, 별도 slide_notes.md 파일로도 저장하세요. PPTX 파일 출력은 절대 생략하지 마세요.\n"
    prefixed_content = SLIDE_GENERATION_PREFIX + prompt_content + export_suffix
    payload = {
        "prompt": prefixed_content,
        "agentProfile": "manus-1.6-max",
        "taskMode": "agent",
        "interactiveMode": False,
        "createShareableLink": True,
        "hide_in_task_list": True,  # M9: Manus UI 오염 방지
    }
    resp = session.post(
        f"{MANUS_API_BASE}/tasks",
        json=payload,
        timeout=(10, 60),  # M4: connect/read 분리
    )

    # C5+C6: HTTP 상태 코드별 처리
    if resp.status_code == 429:
        retry_after = int(resp.headers.get("Retry-After", 60))
        log.warning("[%s] 429 Rate Limit — %d초 후 재시도", filename, retry_after)
        time.sleep(retry_after)
        resp = session.post(
            f"{MANUS_API_BASE}/tasks",
            json=payload,
            timeout=(10, 60),
        )

    if resp.status_code != 200:
        error_detail = resp.text[:300] if resp.text else "No response body"
        raise requests.RequestException(
            f"{resp.status_code} {resp.reason}: {error_detail}"
        )
    data = resp.json()
    task_id = data.get("id") or data.get("task_id") or data.get("taskId")
    share_url = data.get("share_url") or data.get("task_url", "")
    if share_url:
        log.debug("share_url: %s", share_url)
    return task_id, data


def poll_task(session, task_id, filename, max_wait=MAX_WAIT_TIME):
    """Task 완료까지 폴링 (C6: Error 10091 감지 포함)"""
    start = time.time()
    last_status = None
    consecutive_errors = 0

    while time.time() - start < max_wait:
        # C3: Graceful shutdown 체크
        if _shutdown_requested:
            log.warning("[%s] 종료 요청 — 폴링 중단", filename)
            return {"status": "interrupted", "reason": "shutdown_requested"}

        try:
            resp = session.get(
                f"{MANUS_API_BASE}/tasks/{task_id}",
                params={"convert": "true"},
                timeout=(10, 30),  # M4: connect/read 분리
            )
            resp.raise_for_status()
            data = resp.json()
            status = data.get("status", "unknown")
            consecutive_errors = 0  # 성공 시 카운터 리셋

            if status != last_status:
                elapsed = int(time.time() - start)
                log.info("[%s] status=%s (%ds)", filename, status, elapsed)
                last_status = status

            if status in ("completed", "stopped", "done"):
                stop_reason = data.get("stop_reason") or data.get("stopReason", "")
                if status == "stopped" and stop_reason not in ("finish", "completed"):
                    # C6: Error 10091 (High Load) 감지
                    error_code = data.get("error_code") or data.get("errorCode", "")
                    if str(error_code) == "10091" or "high load" in stop_reason.lower():
                        log.warning(
                            "[%s] Error 10091 (High Load Termination) — 재시도 가능",
                            filename,
                        )
                    return {
                        "status": "failed",
                        "reason": stop_reason,
                        "data": data,
                        "error_code": str(error_code),
                    }
                return {"status": "completed", "data": data}

            if status in ("failed", "error", "cancelled"):
                error_code = data.get("error_code") or data.get("errorCode", "")
                return {
                    "status": "failed",
                    "reason": status,
                    "data": data,
                    "error_code": str(error_code),
                }

        except requests.RequestException as e:
            consecutive_errors += 1
            elapsed = int(time.time() - start)
            log.warning(
                "[%s] 폴링 오류 (%ds, 연속 %d회): %s",
                filename,
                elapsed,
                consecutive_errors,
                e,
            )
            # 연속 10회 이상 오류 시 중단
            if consecutive_errors >= 10:
                log.error(
                    "[%s] 연속 폴링 오류 %d회 — 중단", filename, consecutive_errors
                )
                return {
                    "status": "failed",
                    "reason": f"연속 폴링 오류 {consecutive_errors}회",
                }

        time.sleep(POLL_INTERVAL)

    return {"status": "timeout", "reason": f"최대 대기 시간 {max_wait}초 초과"}


def extract_file_urls(task_data):
    """Task 응답에서 다운로드 가능한 파일 URL 추출"""
    urls = []
    seen_urls = set()

    def _add(url, name, mime=""):
        if url and url not in seen_urls:
            seen_urls.add(url)
            urls.append({"url": url, "name": name, "mime": mime})

    output = task_data.get("output") or task_data.get("outputs") or []

    if isinstance(output, list):
        for item in output:
            if not isinstance(item, dict):
                continue
            file_url = item.get("fileUrl") or item.get("file_url") or item.get("url")
            file_name = (
                item.get("fileName") or item.get("file_name") or item.get("name", "")
            )
            mime = item.get("mimeType") or item.get("mime_type", "")
            if file_url:
                _add(file_url, file_name, mime)
            content_list = item.get("content") or []
            if isinstance(content_list, list):
                for c in content_list:
                    if not isinstance(c, dict):
                        continue
                    if c.get("type") == "output_file":
                        c_url = (
                            c.get("fileUrl") or c.get("file_url") or c.get("url", "")
                        )
                        c_name = (
                            c.get("fileName") or c.get("file_name") or c.get("name", "")
                        )
                        c_mime = c.get("mimeType") or c.get("mime_type", "")
                        _add(c_url, c_name, c_mime)

    attachments = task_data.get("attachments") or []
    if isinstance(attachments, list):
        for att in attachments:
            if isinstance(att, dict):
                url = att.get("url") or att.get("fileUrl")
                name = (
                    att.get("file_name") or att.get("fileName") or att.get("name", "")
                )
                _add(url, name, att.get("mimeType", ""))

    artifacts = task_data.get("artifacts") or []
    if isinstance(artifacts, list):
        for art in artifacts:
            if isinstance(art, dict):
                url = art.get("url") or art.get("downloadUrl")
                name = art.get("name") or art.get("fileName", "")
                _add(url, name, art.get("mimeType", ""))

    return urls


def download_file(session, url, save_path):
    """URL에서 파일 다운로드 (C8: 원자적 쓰기 + M5: 최소 크기 검증)"""
    resp = session.get(url, stream=True, timeout=(10, 120))
    resp.raise_for_status()
    save_path.parent.mkdir(parents=True, exist_ok=True)

    # C8: 임시 파일에 먼저 쓰고, 완료 후 rename (원자적 쓰기)
    fd, tmp_path = tempfile.mkstemp(
        dir=str(save_path.parent),
        suffix=".tmp",
    )
    try:
        with os.fdopen(fd, "wb") as f:
            for chunk in resp.iter_content(chunk_size=8192):
                f.write(chunk)
        # rename (원자적 — 실패 시 깨진 파일 방지)
        os.replace(tmp_path, str(save_path))
    except Exception:
        # 실패 시 임시 파일 정리
        try:
            os.unlink(tmp_path)
        except OSError:
            pass
        raise

    file_size = save_path.stat().st_size

    # M5: 최소 파일 크기 검증
    if save_path.suffix.lower() == ".pptx" and file_size < MIN_PPTX_SIZE:
        log.warning(
            "PPTX 파일 크기 의심: %s (%d bytes < %d bytes)",
            save_path.name,
            file_size,
            MIN_PPTX_SIZE,
        )

    return file_size


def strip_headers_footers(pptx_path):
    """PPTX 파일에서 헤더/풋터/페이지번호를 프로그래밍으로 제거 (후처리)"""
    try:
        from pptx import Presentation
        from pptx.util import Emu
        import copy
    except ImportError:
        log.warning("python-pptx 미설치 — 후처리 건너뜀 (pip install python-pptx)")
        return False

    prs = Presentation(str(pptx_path))
    modified = False

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    header_zone = int(slide_h * 0.12)
    footer_zone = int(slide_h * 0.88)

    for slide in prs.slides:
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_idx = shape.placeholder_format.idx
                if ph_idx in (10, 11, 12):
                    shapes_to_remove.append(shape)
                    continue

            if (
                hasattr(shape, "top")
                and hasattr(shape, "height")
                and hasattr(shape, "text")
            ):
                top = shape.top or 0
                height = shape.height or 0
                bottom = top + height
                text = (shape.text or "").strip()
                is_thin = height < int(slide_h * 0.10)
                in_header = top < header_zone and is_thin
                in_footer = bottom > footer_zone and is_thin
                if in_header or in_footer:
                    if len(text) < 80 or not text:
                        shapes_to_remove.append(shape)
                        continue

            name = (shape.name or "").lower()
            if any(
                kw in name
                for kw in (
                    "footer",
                    "header",
                    "slide number",
                    "date placeholder",
                    "page number",
                    "ftr",
                    "hdr",
                    "sldnum",
                )
            ):
                shapes_to_remove.append(shape)

        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
            modified = True

    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            for ph in list(layout.placeholders):
                if ph.placeholder_format.idx in (10, 11, 12):
                    sp = ph._element
                    sp.getparent().remove(sp)
                    modified = True

    if modified:
        prs.save(str(pptx_path))
        log.debug("[후처리] 헤더/풋터/페이지번호 제거 완료")
    else:
        log.debug("[후처리] 제거할 헤더/풋터 없음 (깨끗한 상태)")

    return modified


# ═══════════════════════════════════════════════════════════════
# 분할/병합 함수
# ═══════════════════════════════════════════════════════════════


def _count_slides(content):
    import re

    return len(re.findall(r"^#{2,4}\s*(?:슬라이드|Slide)\s*\d+", content, re.MULTILINE))


def _extract_section(content, section_marker):
    import re

    pattern = rf"^(#{1, 3}\s*{re.escape(section_marker)}.*)$"
    match = re.search(pattern, content, re.MULTILINE)
    if not match:
        return ""
    start = match.start()
    next_section = re.search(
        r"^#{1,3}\s*[①②③④⑤⑥§]", content[match.end() :], re.MULTILINE
    )
    if next_section:
        end = match.end() + next_section.start()
    else:
        end = len(content)
    return content[start:end].rstrip()


def _find_session_boundaries(section_text, pattern):
    import re

    boundaries = []
    for m in re.finditer(pattern, section_text, re.MULTILINE):
        boundaries.append((m.start(), m.group(0).strip()))
    if not boundaries:
        return [(0, len(section_text), "전체")]
    result = []
    for i, (start, label) in enumerate(boundaries):
        end = boundaries[i + 1][0] if i + 1 < len(boundaries) else len(section_text)
        result.append((start, end, label))
    return result


def split_by_session(
    prompt_content, threshold_lines=CHUNK_THRESHOLD, threshold_slides=CHUNK_MAX_SLIDES
):
    """대형 프롬프트를 교시(세션) 경계에서 분할."""
    import re

    lines = prompt_content.splitlines()
    line_count = len(lines)
    slide_count = _count_slides(prompt_content)

    if line_count < threshold_lines and slide_count < threshold_slides:
        return [prompt_content]

    log.info("[분할] %d줄, ~%d슬라이드 → 세션 단위 분할 시작", line_count, slide_count)

    section_3 = _extract_section(prompt_content, "③")
    section_6 = _extract_section(prompt_content, "⑥")

    if not section_3 and not section_6:
        log.info("[분할] ③/⑥ 섹션을 찾을 수 없음 — 분할 건너뜀")
        return [prompt_content]

    s3_boundaries = (
        _find_session_boundaries(section_3, r"^####\s*\d+\.\s*세션")
        if section_3
        else []
    )

    s6_boundaries = (
        _find_session_boundaries(section_6, r"^##\s*세션") if section_6 else []
    )

    split_count = max(len(s3_boundaries), len(s6_boundaries))
    if split_count <= 1:
        log.info("[분할] 세션 경계가 1개 이하 — 분할 건너뜀")
        return [prompt_content]

    common_header_parts = []
    for marker in ["①", "②", "④", "⑤"]:
        section = _extract_section(prompt_content, marker)
        if section:
            common_header_parts.append(section)
    common_header = "\n\n".join(common_header_parts)

    section_3_header = ""
    if section_3:
        first_boundary = re.search(r"^####\s*\d+\.\s*세션", section_3, re.MULTILINE)
        if first_boundary and first_boundary.start() > 0:
            section_3_header = section_3[: first_boundary.start()].rstrip()

    section_6_header = ""
    if section_6:
        first_boundary = re.search(r"^##\s*세션", section_6, re.MULTILINE)
        if first_boundary and first_boundary.start() > 0:
            section_6_header = section_6[: first_boundary.start()].rstrip()

    chunks = []
    for i in range(split_count):
        parts = [common_header]
        if i < len(s3_boundaries):
            start, end, label = s3_boundaries[i]
            s3_chunk = section_3[start:end].rstrip()
            if section_3_header:
                parts.append(section_3_header + "\n\n" + s3_chunk)
            else:
                parts.append(s3_chunk)
        if i < len(s6_boundaries):
            start, end, label = s6_boundaries[i]
            s6_chunk = section_6[start:end].rstrip()
            if section_6_header:
                parts.append(section_6_header + "\n\n" + s6_chunk)
            else:
                parts.append(s6_chunk)
        chunk_text = "\n\n".join(parts)
        chunks.append(chunk_text)

    log.info("[분할] %d개 청크로 분할 완료", split_count)
    for i, c in enumerate(chunks):
        log.debug("  청크 %d: %d줄", i + 1, len(c.splitlines()))

    return chunks


def merge_pptx(chunk_pptx_paths, output_path):
    """분할 제출된 청크 PPTX들을 하나의 파일로 병합."""
    try:
        from pptx import Presentation
        import copy
    except ImportError:
        log.warning("python-pptx 미설치 — 병합 건너뜀")
        return 0

    valid_paths = [p for p in chunk_pptx_paths if Path(p).exists()]
    if not valid_paths:
        log.warning("병합할 PPTX 파일 없음")
        return 0

    if len(valid_paths) == 1:
        shutil.copy2(valid_paths[0], output_path)
        prs = Presentation(str(valid_paths[0]))
        return len(prs.slides)

    base_prs = Presentation(str(valid_paths[0]))
    total_slides = len(base_prs.slides)

    for pptx_path in valid_paths[1:]:
        src_prs = Presentation(str(pptx_path))
        for slide in src_prs.slides:
            slide_layout = base_prs.slide_layouts[6]
            for layout in base_prs.slide_layouts:
                if layout.name == "Blank" or "blank" in layout.name.lower():
                    slide_layout = layout
                    break

            new_slide = base_prs.slides.add_slide(slide_layout)

            for elem in list(new_slide.shapes._spTree):
                if elem.tag.endswith("}sp") or elem.tag.endswith("}pic"):
                    new_slide.shapes._spTree.remove(elem)

            for shape in slide.shapes:
                el = copy.deepcopy(shape._element)
                new_slide.shapes._spTree.append(el)

            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text:
                    if not new_slide.has_notes_slide:
                        new_slide.notes_slide
                    new_slide.notes_slide.notes_text_frame.text = notes_text

            total_slides += 1

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    base_prs.save(str(output_path))
    log.info("[병합] %d개 PPTX → %d슬라이드 병합 완료", len(valid_paths), total_slides)
    return total_slides


# ═══════════════════════════════════════════════════════════════
# 처리 함수
# ═══════════════════════════════════════════════════════════════


def _submit_single_content(session, prompt_content, label, output_dir):
    """단일 콘텐츠 제출→폴링→다운로드 (C6: 10091 자동 재제출 포함)"""
    max_submit_attempts = 2  # 10091 시 1회 재제출

    for attempt in range(1, max_submit_attempts + 1):
        try:
            task_id, create_data = create_task(session, prompt_content, label)
            log.info("task_id: %s (시도 %d/%d)", task_id, attempt, max_submit_attempts)
        except requests.RequestException as e:
            log.error("[%s] Task 생성 실패 (시도 %d): %s", label, attempt, e)
            if attempt < max_submit_attempts:
                wait = RETRY_BACKOFF * (2 ** (attempt - 1))
                log.info("[%s] %d초 후 재시도...", label, wait)
                time.sleep(wait)
                continue
            return {"label": label, "status": "submit_failed", "error": str(e)}

        result = poll_task(session, task_id, label)

        # C6: Error 10091 감지 → 대기 후 재제출
        if (
            result["status"] == "failed"
            and result.get("error_code") == "10091"
            and attempt < max_submit_attempts
        ):
            log.warning("[%s] Error 10091 — 60초 대기 후 재제출", label)
            time.sleep(60)
            continue

        if result["status"] != "completed":
            log.error("[%s] %s", label, result.get("reason", "unknown"))
            return {
                "label": label,
                "task_id": task_id,
                "status": result["status"],
                "reason": result.get("reason", ""),
            }

        # 성공 — 파일 다운로드
        task_data = result["data"]
        file_urls = extract_file_urls(task_data)

        if not file_urls:
            log.warning("[%s] 다운로드 가능한 파일 없음", label)
            share_link = task_data.get("shareableLink") or task_data.get(
                "shareable_link", ""
            )
            return {
                "label": label,
                "task_id": task_id,
                "status": "completed_no_file",
                "shareable_link": share_link,
            }

        downloaded = []
        for finfo in file_urls:
            raw_name = finfo.get("name", "")
            ext = Path(raw_name).suffix if raw_name else ""
            if not ext:
                url_path = finfo["url"].split("?")[0]
                ext = Path(url_path).suffix if "." in Path(url_path).name else ".pptx"

            if ext.lower() == ".pptx":
                save_name = f"{label}.pptx"
            else:
                base = Path(raw_name).stem if raw_name else f"{label}_notes"
                save_name = f"{label}_{base}{ext}" if raw_name else f"{label}{ext}"

            save_path = output_dir / save_name
            try:
                size = download_file(session, finfo["url"], save_path)
                log.info("[OK] 다운로드: %s (%s bytes)", save_name, f"{size:,}")
                if ext.lower() == ".pptx":
                    strip_headers_footers(save_path)
                    size = save_path.stat().st_size
                downloaded.append({"path": str(save_path), "size": size})
            except (requests.RequestException, OSError) as e:
                log.error("다운로드 실패 (%s): %s", save_name, e)
                downloaded.append({"url": finfo["url"], "error": str(e)})

        return {
            "label": label,
            "task_id": task_id,
            "status": "completed",
            "downloaded": downloaded,
        }

    # 모든 시도 실패
    return {"label": label, "status": "submit_failed", "error": "최대 재시도 초과"}


def process_single(session, prompt_file, output_dir, no_split=False):
    """단일 프롬프트 파일 처리: (분할 판단→) 제출 → 폴링 → 다운로드 (→ 병합)"""
    filename = prompt_file.stem
    short_name = filename.replace("_슬라이드 생성 프롬프트", "")
    log.info("%s", "=" * 60)
    log.info("[제출] %s", short_name)

    prompt_content = prompt_file.read_text(encoding="utf-8")

    if no_split:
        chunks = [prompt_content]
    else:
        chunks = split_by_session(prompt_content)

    if len(chunks) == 1:
        result = _submit_single_content(session, chunks[0], short_name, output_dir)
        return {
            "file": str(prompt_file.name),
            "task_id": result.get("task_id"),
            "status": result["status"],
            "chunks": 1,
            "downloaded": result.get("downloaded", []),
            "shareable_link": result.get("shareable_link", ""),
            "reason": result.get("reason", ""),
            "error": result.get("error", ""),
        }

    log.info("[%s] %d개 청크 순차 제출", short_name, len(chunks))
    chunk_results = []
    chunk_pptx_paths = []
    chunk_dir = output_dir / f".chunks_{short_name}"
    chunk_dir.mkdir(exist_ok=True)

    # C4: 청크 단위 중간 상태 저장
    chunk_log_path = chunk_dir / "chunk_progress.json"

    for i, chunk_content in enumerate(chunks):
        # C3: 종료 요청 체크
        if _shutdown_requested:
            log.warning(
                "[%s] 종료 요청 — 청크 %d/%d 에서 중단", short_name, i + 1, len(chunks)
            )
            break

        chunk_label = f"{short_name}_chunk{i + 1}"
        log.info("--- 청크 %d/%d ---", i + 1, len(chunks))
        cr = _submit_single_content(session, chunk_content, chunk_label, chunk_dir)
        chunk_results.append(cr)

        # C4: 청크 진행 상태 즉시 저장
        chunk_log_path.write_text(
            json.dumps(chunk_results, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )

        for dl in cr.get("downloaded", []):
            if dl.get("path", "").endswith(".pptx"):
                chunk_pptx_paths.append(dl["path"])

    all_completed = all(cr["status"] == "completed" for cr in chunk_results)

    if chunk_pptx_paths and len(chunk_pptx_paths) > 1:
        merged_path = output_dir / f"{short_name}.pptx"
        total_slides = merge_pptx(chunk_pptx_paths, str(merged_path))
        if total_slides > 0:
            strip_headers_footers(merged_path)
            merged_size = merged_path.stat().st_size
            downloaded = [{"path": str(merged_path), "size": merged_size}]
        else:
            downloaded = [
                {"path": p, "size": Path(p).stat().st_size}
                for p in chunk_pptx_paths
                if Path(p).exists()
            ]
    elif chunk_pptx_paths:
        single_path = output_dir / f"{short_name}.pptx"
        shutil.copy2(chunk_pptx_paths[0], single_path)
        downloaded = [{"path": str(single_path), "size": single_path.stat().st_size}]
    else:
        downloaded = []

    return {
        "file": str(prompt_file.name),
        "task_id": ",".join(cr.get("task_id", "") or "" for cr in chunk_results),
        "status": "completed" if all_completed else "partial",
        "chunks": len(chunks),
        "chunk_results": chunk_results,
        "downloaded": downloaded,
    }


def filter_prompts(prompts, file_filter):
    """--file 옵션으로 프롬프트 목록을 필터링"""
    if not file_filter:
        return prompts

    matched = []
    for keyword in file_filter:
        kw_lower = keyword.lower()
        for p in prompts:
            name_lower = p.name.lower()
            stem_lower = p.stem.lower()
            if (
                kw_lower == name_lower
                or kw_lower == stem_lower
                or kw_lower in name_lower
                or kw_lower in stem_lower
            ):
                if p not in matched:
                    matched.append(p)

    return matched


# ═══════════════════════════════════════════════════════════════
# 메인
# ═══════════════════════════════════════════════════════════════


def main():
    global _current_task_log, _current_output_dir

    parser = argparse.ArgumentParser(
        description="Manus AI 슬라이드 생성",
        epilog="예시: python .agent/scripts/manus_slide.py project/ --file Day1_AM",
    )
    parser.add_argument("project_dir", nargs="?", help="프로젝트 폴더 경로")
    parser.add_argument(
        "--file",
        "-f",
        nargs="+",
        dest="file_filter",
        help="처리할 파일 필터 (부분 매칭 지원, 복수 지정 가능)",
    )
    parser.add_argument(
        "--resume", action="store_true", help="이전 완료 파일을 자동 스킵하여 이어하기"
    )
    parser.add_argument(
        "--dry-run", action="store_true", help="실제 API 호출 없이 테스트"
    )
    parser.add_argument(
        "--no-split",
        action="store_true",
        help="세션 단위 분할 비활성화 (원본 전체 제출)",
    )
    parser.add_argument(
        "--quiet", "-q", action="store_true", help="사일런스 모드 (경고/에러만 출력)"
    )
    parser.add_argument(
        "--verbose", "-v", action="store_true", help="디버그 모드 (상세 출력)"
    )
    args = parser.parse_args()

    # M1+M2: 로깅 설정
    _setup_logging(quiet=args.quiet, verbose=args.verbose)

    # C3: 시그널 핸들러 등록
    signal.signal(signal.SIGINT, _handle_signal)
    signal.signal(signal.SIGTERM, _handle_signal)
    atexit.register(_save_checkpoint)

    # API Key 확인
    api_key = get_api_key()
    if not api_key:
        log.error("MANUS_API_KEY를 찾을 수 없습니다. .agent/.env 파일에 설정하세요.")
        sys.exit(1)
    log.info("API Key 확인됨: %s...%s", api_key[:10], api_key[-4:])

    # C7: requests.Session 생성
    session = _create_session(api_key)

    # 프로젝트 폴더 탐색
    project = find_project_folder(args.project_dir)
    if not project:
        log.error("04_SlidePrompt/ 폴더를 찾을 수 없습니다. 경로를 인자로 지정하세요.")
        sys.exit(1)
    log.info("프로젝트: %s", project.name)

    # 프롬프트 파일 탐색
    prompts = discover_prompts(project)
    if not prompts:
        log.error("%s 에 프롬프트 파일이 없습니다.", project / "04_SlidePrompt")
        sys.exit(1)

    # --file 필터 적용
    if args.file_filter:
        prompts = filter_prompts(prompts, args.file_filter)
        if not prompts:
            log.error(
                "--file 필터 '%s'에 매칭되는 파일이 없습니다.",
                " ".join(args.file_filter),
            )
            sys.exit(1)
        log.info("--file 필터 적용: %d개 선택", len(prompts))
    log.info("프롬프트 파일 %d개 발견", len(prompts))
    for p in prompts:
        log.info("  %s", p.name)

    # 출력 디렉토리
    output_dir = project / "07_ManusSlides"
    output_dir.mkdir(exist_ok=True)
    _current_output_dir = output_dir
    log.info("출력 폴더: %s", output_dir)

    # M8: Lock File — 중복 실행 방지
    if not _acquire_lock(output_dir):
        sys.exit(1)
    atexit.register(_release_lock, output_dir)

    # M6: 디스크 공간 확인
    if not _check_disk_space(output_dir):
        _release_lock(output_dir)
        sys.exit(1)

    if args.dry_run:
        log.info("[DRY-RUN] 실제 API 호출을 건너뜁니다.")
        _release_lock(output_dir)
        return

    # M3: Preflight 연결 테스트
    if not _preflight_check(session):
        log.error("Manus API 연결 실패 — 실행을 중단합니다.")
        _release_lock(output_dir)
        sys.exit(1)

    # C2: --resume — 이전 완료 파일 스킵
    completed_files = set()
    if args.resume:
        completed_files = _load_completed_files(output_dir)

    # 실행 시작
    start_time = time.time()
    log.info("%s", "=" * 60)
    log.info("슬라이드 생성 시작 — %s", datetime.now().strftime("%H:%M:%S"))
    log.info("%s", "=" * 60)

    results = []
    _current_task_log = []

    # 순차 처리 (Manus API 제한 고려)
    for prompt_file in prompts:
        # C3: Graceful shutdown 체크
        if _shutdown_requested:
            log.warning("종료 요청 — 남은 파일 %d개 스킵", len(prompts) - len(results))
            break

        # C2: 이전 완료 파일 스킵
        if prompt_file.name in completed_files:
            log.info("[SKIP] %s (이전 완료)", prompt_file.name)
            continue

        result = process_single(
            session, prompt_file, output_dir, no_split=args.no_split
        )
        results.append(result)

        # 중간 로그 저장 (중단 복구용)
        _current_task_log.append(result)
        _save_checkpoint()

    # M7: 임시 청크 디렉토리 정리
    _cleanup_chunk_dirs(output_dir)

    # 최종 리포트
    elapsed = int(time.time() - start_time)
    log.info("%s", "=" * 60)
    log.info("슬라이드 생성 완료 — 총 %d초 소요", elapsed)
    log.info("%s", "=" * 60)

    success = [r for r in results if r["status"] == "completed"]
    no_file = [r for r in results if r["status"] == "completed_no_file"]
    failed = [
        r for r in results if r["status"] not in ("completed", "completed_no_file")
    ]

    log.info("성공 (PPTX 다운로드): %d개", len(success))
    for r in success:
        for d in r.get("downloaded", []):
            if "path" in d:
                log.info("  %s (%s bytes)", Path(d["path"]).name, f"{d['size']:,}")

    if no_file:
        log.info("완료 (수동 다운로드 필요): %d개", len(no_file))
        for r in no_file:
            link = r.get("shareable_link", "N/A")
            log.info("  %s → %s", r["file"], link)

    if failed:
        log.warning("실패: %d개", len(failed))
        for r in failed:
            log.warning(
                "  %s: %s", r["file"], r.get("reason", r.get("error", "unknown"))
            )

    log.info("로그: %s", output_dir / "manus_task_log.json")

    # 최종 리포트 저장
    report = {
        "generated_at": datetime.now().isoformat(),
        "project": str(project),
        "total_files": len(prompts),
        "skipped_resume": len(completed_files),
        "success": len(success),
        "no_file": len(no_file),
        "failed": len(failed),
        "interrupted": _shutdown_requested,
        "elapsed_seconds": elapsed,
        "results": results,
    }
    report_path = output_dir / "generation_report.json"
    report_path.write_text(
        json.dumps(report, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    log.info("리포트: %s", report_path)


if __name__ == "__main__":
    main()
